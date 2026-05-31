"""Comprehensive SEO + GEO strategy deck builder.

Produces a 30-40 slide strategic deck combining technical audit data with
Claude-generated narrative analysis. Each slide is its own function for
maintainability.
"""
from __future__ import annotations

import io
from datetime import datetime
from typing import Iterable

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
LIGHT2 = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- HELPERS ----------
def _add_title(slide, text: str, subtitle: str = "", color=NAVY):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY


def _accent_bar(slide, color=TEAL, y=1.18, w=1.2):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(w), Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _footer(slide, text: str):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY


def _score_color(score: float) -> RGBColor:
    if score >= 75:
        return GREEN
    if score >= 50:
        return AMBER
    return RED


def _kpi_card(slide, x, y, w, h, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = LIGHT
    tb = slide.shapes.add_textbox(x, y + Inches(0.15), w, h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(value)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY


def _add_paragraphs(slide, x, y, w, h, items: list[tuple[str, dict] | str], title: str = "", title_color=NAVY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = title_color
        first = False
    for item in items:
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        if not text:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(opts.get("size", 11))
        p.font.bold = opts.get("bold", False)
        p.font.italic = opts.get("italic", False)
        p.font.color.rgb = opts.get("color", NAVY)
        if "space_before" in opts:
            p.space_before = Pt(opts["space_before"])


def _bullets(slide, x, y, w, h, items: Iterable[str], title: str = "", title_color=NAVY, size=11):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = title_color
        first = False
    for item in items:
        if not item:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY


def _table(slide, x, y, w, h, headers: list[str], rows: list[list[str]],
            col_widths: list[float] | None = None, status_col: int | None = None,
            status_pass: str = "ALLOWED"):
    """Generic table helper. col_widths in inches."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            if i < n_cols:
                table.columns[i].width = Inches(cw)
    for i, h_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = NAVY
                    if status_col is not None and c_idx == status_col:
                        run.font.bold = True
                        run.font.color.rgb = GREEN if str(val).upper() == status_pass.upper() else RED
    return table


def _bar_chart_image(labels: list[str], values: list[float], colors: list[str], title: str) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(8, 4.2), dpi=140)
    bars = ax.barh(labels, values, color=colors)
    ax.set_xlim(0, 100)
    ax.invert_yaxis()
    ax.set_xlabel("% pass", fontsize=10)
    ax.set_title(title, fontsize=12, weight="bold", color="#0B1F3A", loc="left")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(labelsize=9)
    for bar, v in zip(bars, values):
        ax.text(min(v + 1.5, 96), bar.get_y() + bar.get_height() / 2,
                f"{int(round(v))}%", va="center", fontsize=9, color="#0B1F3A")
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _score_to_color_hex(score: float) -> str:
    if score >= 75:
        return "#16A34A"
    if score >= 50:
        return "#F59E0B"
    return "#DC2626"


# ---------- SLIDE FUNCTIONS ----------
def slide_cover(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.1), prs.slide_width, Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(12), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "SEO + GEO Strategy & Audit"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(12), Inches(0.7))
    p = tb2.text_frame.paragraphs[0]
    p.text = ctx["domain"]
    p.font.size = Pt(28)
    p.font.color.rgb = TEAL

    tb3 = s.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(12), Inches(0.6))
    p = tb3.text_frame.paragraphs[0]
    p.text = f"Generated {datetime.now().strftime('%B %d, %Y')}  •  {ctx['pages_analyzed']} pages analyzed  •  {len(ctx['competitors_metrics'])} competitors"
    p.font.size = Pt(15)
    p.font.color.rgb = LIGHT

    tb4 = s.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(12), Inches(0.5))
    p = tb4.text_frame.paragraphs[0]
    p.text = "Traditional SEO  •  Generative Engine Optimization  •  Competitor Analysis  •  Content Strategy  •  Backlinks  •  GMB"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT
    p.font.italic = True


def slide_scope(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Scope of This Document", "Purpose & objectives")
    _accent_bar(s)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ctx["insights"].get("scope", "") or (
        f"The purpose of this document is to outline the SEO and Generative Engine Optimization "
        f"strategy for {ctx['domain']}, with the goal of improving online visibility, driving organic "
        f"traffic, and boosting conversions. This strategy combines technical SEO health, content "
        f"strategy, competitor analysis, and AI-search readiness to deliver measurable growth."
    )
    p.font.size = Pt(15)
    p.font.color.rgb = NAVY


def slide_toc(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Table of Contents")
    _accent_bar(s)
    sections = [
        "1.  Scope of Document",
        "2.  Domain Rating & Authority — Comparison + Analysis",
        "3.  Traffic, Keywords & Traffic Value — Comparison + Analysis",
        "4.  Keyword Research & Page-Level Mapping",
        "5.  Recommended Website Navigation Layout",
        "6.  Content Strategy — Topic Silos & Blog Calendar",
        "7.  Google My Business Audit & Recommendations",
        "8.  Technical SEO Audit — Per-Page Findings",
        "9.  Core Web Vitals — Sampled Performance",
        "10. Generative Engine Optimization (GEO) — Detailed",
        "11. Priority Recommendations & Next Steps",
    ]
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(sections):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = NAVY


def slide_dr_explainer(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Domain Rating & Domain Authority", "What these metrics measure")
    _accent_bar(s)

    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True
    paras = [
        ("Domain Rating (DR)", {"size": 17, "bold": True, "color": TEAL}),
        ("Ahrefs' 0–100 metric measuring the strength of a domain's backlink profile. "
         "Higher DR correlates with the ability to rank for competitive keywords. "
         "DR is logarithmic — moving from DR 20 to 30 is significantly harder than moving from 10 to 20.",
         {"size": 12}),
        ("Domain Authority (DA)", {"size": 17, "bold": True, "color": TEAL, "space_before": 12}),
        ("Moz's 0–100 metric that predicts how well a website will rank on search engines. "
         "DA evaluates linking root domains and total number of links into a single score. "
         "Like DR, DA is logarithmic and most useful as a comparative measure.",
         {"size": 12}),
        ("Why they matter together", {"size": 17, "bold": True, "color": TEAL, "space_before": 12}),
        ("DR and DA are the two most-cited authority metrics in industry tooling. "
         "Tracking both gives a balanced view of authority and helps benchmark against competitors. "
         "They are not used directly by Google but correlate strongly with ranking ability.",
         {"size": 12}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), paras)


def slide_dr_table(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Domain Rating & Authority — Comparison", "Your site vs. competitors")
    _accent_bar(s)

    def _dr_str(v):
        try:
            return str(int(v)) if int(v) > 0 else "—"
        except (TypeError, ValueError):
            return "—"
    rows = [
        [ctx["your_site_metrics"]["domain"],
         _dr_str(ctx["your_site_metrics"].get("dr")),
         _dr_str(ctx["your_site_metrics"].get("da")),
         "★  Your site"],
    ]
    for c in ctx["competitors_metrics"]:
        rows.append([
            c["domain"],
            _dr_str(c.get("dr")),
            _dr_str(c.get("da")),
            "Competitor",
        ])
    _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4 + 0.5 * (len(rows) + 1)),
           headers=["Domain", "Domain Rating (DR)", "Domain Authority (DA)", "Role"],
           rows=rows, col_widths=[5.5, 2.3, 2.3, 2.2])
    # Note about data source when synthetic
    notes = ctx["your_site_metrics"].get("_source_notes", []) or []
    if "authority-synth-from-organic" in notes:
        _footer(s, "DR/DA estimated from organic-traffic metrics (DataForSEO Labs). For exact Ahrefs DR + Moz DA, paste them in the input form or activate DataForSEO Backlinks API.")


def slide_dr_analysis(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Domain Authority — Analysis", "Where we stand and how to close the gap")
    _accent_bar(s)
    ins = ctx["insights"].get("authority", {})

    paras = [
        ("Narrative", {"size": 14, "bold": True, "color": TEAL}),
        (ins.get("narrative", "(analysis pending)"), {"size": 11}),
        ("Objective", {"size": 14, "bold": True, "color": TEAL, "space_before": 10}),
        (ins.get("objective", ""), {"size": 11}),
        ("Reason", {"size": 14, "bold": True, "color": TEAL, "space_before": 10}),
        (ins.get("reason", ""), {"size": 11}),
        ("How to achieve it", {"size": 14, "bold": True, "color": TEAL, "space_before": 10}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.3), paras)
    _bullets(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
             ins.get("how_to_achieve", []), title="", size=11)


def slide_traffic_table(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Traffic, Keywords & Traffic Value — Comparison", "All values from Ahrefs / DataForSEO")
    _accent_bar(s)

    rows = [
        [ctx["your_site_metrics"]["domain"],
         _fmt_int(ctx['your_site_metrics'].get('ranking_keywords', 0)),
         _fmt_int(ctx['your_site_metrics'].get('organic_traffic', 0)),
         _fmt_money(ctx['your_site_metrics'].get('traffic_value', 0)),
         "★  Your site"],
    ]
    for c in ctx["competitors_metrics"]:
        rows.append([
            c["domain"],
            _fmt_int(c.get('ranking_keywords', 0)),
            _fmt_int(c.get('organic_traffic', 0)),
            _fmt_money(c.get('traffic_value', 0)),
            "Competitor",
        ])
    _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4 + 0.5 * (len(rows) + 1)),
           headers=["Domain", "Ranking Keywords", "Organic Traffic", "Traffic Value", "Role"],
           rows=rows, col_widths=[4.2, 2.2, 2.2, 2.2, 1.5])


def slide_traffic_analysis(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Traffic Analysis", "Facts, traffic-source mix, and top pages")
    _accent_bar(s)
    ins = ctx["insights"].get("traffic", {})

    paras = [
        ("Facts", {"size": 14, "bold": True, "color": TEAL}),
        (ins.get("our_facts", ""), {"size": 11}),
        ("Traffic Source Analysis", {"size": 14, "bold": True, "color": TEAL, "space_before": 10}),
        (ins.get("traffic_source_analysis", ""), {"size": 11}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(3), paras)

    # Top pages table
    top_pages = ins.get("our_top_pages", []) or []
    if top_pages:
        _bullets(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3),
                 [f"{p.get('url', '')}  —  ({p.get('keyword', '')})  →  {p.get('insight', '')}" for p in top_pages[:4]],
                 title="Our top traffic-driving pages", title_color=TEAL, size=10)


def slide_competitor_traffic_analysis(prs, ctx):
    """One slide per competitor (or up to N competitors per slide if many)."""
    blank = prs.slide_layouts[6]
    ins = ctx["insights"].get("traffic", {})
    summaries = ins.get("competitor_summaries", []) or []
    for comp in summaries[:5]:
        s = prs.slides.add_slide(blank)
        _add_title(s, f"Competitor Deep-Dive: {comp.get('domain', '')}",
                   "Where their organic traffic comes from")
        _accent_bar(s, color=AMBER)

        _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2),
                        [(comp.get("insight", ""), {"size": 12})])
        pages = comp.get("top_pages", []) or []
        if pages:
            _bullets(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.3),
                     [f"{p.get('url', '')}  —  Keyword: {p.get('keyword', '')}" for p in pages[:6]],
                     title="Key traffic-driving pages", title_color=TEAL, size=10)


def slide_keyword_homepage(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Keyword Research — Homepage", "Primary + supporting keywords for the home page")
    _accent_bar(s)
    km = ctx["insights"].get("keyword_map", {})
    if not isinstance(km, dict):
        km = {}
    home = km.get("homepage", {})
    if not isinstance(home, dict):
        home = {}
    primary = str(home.get("primary", "") or "")
    secondary = [str(x) for x in (home.get("secondary", []) or []) if x]

    _add_paragraphs(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.5),
                    [("Primary keyword", {"size": 14, "bold": True, "color": TEAL}),
                     (primary or "—", {"size": 22, "bold": True})])
    _bullets(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(3.3),
             secondary, title="Supporting keywords", title_color=TEAL, size=13)


def slide_keyword_services(prs, ctx):
    """Service pages keyword mapping. Splits across slides if many services."""
    blank = prs.slide_layouts[6]
    km = ctx["insights"].get("keyword_map", {})
    if not isinstance(km, dict):
        km = {}
    services = [s for s in (km.get("service_pages", []) or []) if isinstance(s, dict)]
    if not services:
        return
    chunks = [services[i:i + 3] for i in range(0, len(services), 3)]
    for idx, chunk in enumerate(chunks):
        s = prs.slides.add_slide(blank)
        _add_title(s, f"Keyword Research — Service Pages ({idx + 1}/{len(chunks)})",
                   "Page-level keyword mapping with sub-services")
        _accent_bar(s)

        # 3 columns
        col_w = (12.3 - 0.6) / 3
        for c_i, svc in enumerate(chunk):
            x = Inches(0.5 + c_i * (col_w + 0.3))
            # service title box
            tb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(col_w), Inches(0.6))
            tb.fill.solid()
            tb.fill.fore_color.rgb = TEAL
            tb.line.fill.background()
            txt = s.shapes.add_textbox(x, Inches(1.55), Inches(col_w), Inches(0.5))
            p = txt.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = svc.get("page_name", "Service")
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE

            # Content
            content = [
                (f"URL: {svc.get('url_slug', '')}", {"size": 9, "italic": True, "color": GRAY}),
                ("Primary:", {"size": 10, "bold": True, "color": NAVY, "space_before": 6}),
                (svc.get("primary", ""), {"size": 11, "bold": True}),
                ("Secondary:", {"size": 10, "bold": True, "color": NAVY, "space_before": 6}),
            ]
            for sk in (svc.get("secondary", []) or [])[:5]:
                content.append((f"  • {sk}", {"size": 9}))
            subs = [s for s in (svc.get("sub_services", []) or []) if isinstance(s, dict)]
            if subs:
                content.append(("Sub-services:", {"size": 10, "bold": True, "color": NAVY, "space_before": 6}))
                for sub in subs[:4]:
                    content.append((f"  ▸ {sub.get('name', '')} → {sub.get('primary', '')}", {"size": 9}))
            _add_paragraphs(s, x, Inches(2.2), Inches(col_w), Inches(4.7), content)


def slide_keyword_additional(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Keyword Research — Other Key Pages", "About, Contact & additional pages")
    _accent_bar(s)
    km = ctx["insights"].get("keyword_map", {})
    if not isinstance(km, dict):
        km = {}

    pages = []
    about = km.get("about_page")
    if isinstance(about, dict):
        pages.append({"page_name": "About", "url_slug": about.get("url_slug", "/about/"),
                      "primary": about.get("primary", ""), "secondary": about.get("secondary", []) or []})
    contact = km.get("contact_page")
    if isinstance(contact, dict):
        pages.append({"page_name": "Contact", "url_slug": contact.get("url_slug", "/contact/"),
                      "primary": contact.get("primary", ""), "secondary": contact.get("secondary", []) or []})
    for ap in km.get("additional_pages", []) or []:
        if isinstance(ap, dict):
            pages.append(ap)

    if not pages:
        return

    col_w = (12.3 - 0.6) / max(1, min(len(pages), 3))
    for c_i, pg in enumerate(pages[:3]):
        x = Inches(0.5 + c_i * (col_w + 0.3))
        tb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(col_w), Inches(0.6))
        tb.fill.solid()
        tb.fill.fore_color.rgb = AMBER
        tb.line.fill.background()
        txt = s.shapes.add_textbox(x, Inches(1.65), Inches(col_w), Inches(0.5))
        p = txt.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = pg.get("page_name", "Page")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE

        content = [
            (f"URL: {pg.get('url_slug', '')}", {"size": 10, "italic": True, "color": GRAY}),
            ("Primary:", {"size": 11, "bold": True, "color": NAVY, "space_before": 8}),
            (pg.get("primary", ""), {"size": 13, "bold": True}),
            ("Secondary:", {"size": 11, "bold": True, "color": NAVY, "space_before": 8}),
        ]
        for sk in (pg.get("secondary", []) or [])[:6]:
            content.append((f"  • {sk}", {"size": 10}))
        _add_paragraphs(s, x, Inches(2.3), Inches(col_w), Inches(4.6), content)


def slide_navigation(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Recommended Website Navigation", "Information architecture for primary nav")
    _accent_bar(s)
    nav = ctx["insights"].get("navigation", {})
    if not isinstance(nav, dict):
        nav = {}
    items = [n for n in (nav.get("nav_items", []) or []) if isinstance(n, dict)]
    rationale = str(nav.get("rationale", "") or "")

    # Draw nav like a tree
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for nav_item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(6)
        p.text = f"▸ {nav_item.get('label', '')}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = TEAL
        for child in (nav_item.get("children", []) or []):
            if not isinstance(child, dict):
                continue
            p2 = tf.add_paragraph()
            p2.text = f"    ◦ {child.get('label', '')}"
            p2.font.size = Pt(11)
            p2.font.color.rgb = NAVY
            for grandchild in (child.get("children", []) or []):
                if not isinstance(grandchild, dict):
                    continue
                p3 = tf.add_paragraph()
                p3.text = f"        – {grandchild.get('label', '')}"
                p3.font.size = Pt(10)
                p3.font.color.rgb = GRAY

    if rationale:
        _add_paragraphs(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
                        [("Rationale: " + rationale, {"size": 10, "italic": True, "color": GRAY})])


def slide_content_strategy_intro(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Content Strategy — Topic Silos", "Middle-of-funnel blog ideas grouped by silo")
    _accent_bar(s)
    silos_raw = ctx["insights"].get("silos", {})
    if not isinstance(silos_raw, dict):
        silos_raw = {}
    silos = [s for s in (silos_raw.get("silos", []) or []) if isinstance(s, dict)]

    paras = [
        ("Approach", {"size": 14, "bold": True, "color": TEAL}),
        (f"We propose {len(silos)} topic silos, each with 10 middle-of-funnel blog ideas. "
         "MOFU content (comparisons, how-tos, decision guides) drives qualified traffic and "
         "moves prospects from awareness to consideration.",
         {"size": 12}),
        ("Why silos?", {"size": 14, "bold": True, "color": TEAL, "space_before": 12}),
        ("Topic silos signal subject-matter authority to search engines and AI engines alike. "
         "Internal linking within a silo reinforces topical relevance, and each pillar keyword "
         "anchors a cluster that earns ranking power collectively.",
         {"size": 12}),
        ("Why MOFU?", {"size": 14, "bold": True, "color": TEAL, "space_before": 12}),
        ("Top-of-funnel content attracts volume but few buyers. Bottom-of-funnel content converts "
         "but reaches few prospects. MOFU is the sweet spot — high commercial intent, manageable competition.",
         {"size": 12}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.3), paras)


def slide_silo(prs, silo: dict, idx: int, total: int):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, f"Silo {idx}/{total}: {silo.get('silo_name', '')}",
               f"Pillar: {silo.get('pillar_keyword', '')}  |  Audience: {silo.get('audience', '')}")
    _accent_bar(s, color=TEAL)

    blogs = silo.get("blogs", []) or []
    rows = []
    for i, blog in enumerate(blogs[:10], start=1):
        rows.append([
            str(i),
            blog.get("title", ""),
            blog.get("primary_keyword", ""),
            (blog.get("intent", "") or "").title(),
        ])
    if rows:
        _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4 + 0.5 * (len(rows) + 1)),
               headers=["#", "Blog Title", "Primary Keyword", "Intent"],
               rows=rows, col_widths=[0.5, 7.0, 3.0, 1.8])


def _fmt_int(v) -> str:
    try:
        n = int(v or 0)
    except (TypeError, ValueError):
        return "—"
    return f"{n:,}" if n > 0 else "—"


def _fmt_money(v) -> str:
    try:
        n = int(v or 0)
    except (TypeError, ValueError):
        return "—"
    return f"${n:,}" if n > 0 else "—"


def slide_backlinks_table(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Backlinks & Referring Domains — Comparison", "Authority foundation comparison")
    _accent_bar(s)

    rows = [
        [ctx["your_site_metrics"]["domain"],
         _fmt_int(ctx['your_site_metrics'].get('backlinks', 0)),
         _fmt_int(ctx['your_site_metrics'].get('referring_domains', 0)),
         "★  Your site"],
    ]
    for c in ctx["competitors_metrics"]:
        rows.append([
            c["domain"],
            _fmt_int(c.get('backlinks', 0)),
            _fmt_int(c.get('referring_domains', 0)),
            "Competitor",
        ])
    _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4 + 0.5 * (len(rows) + 1)),
           headers=["Domain", "Total Backlinks", "Referring Domains", "Role"],
           rows=rows, col_widths=[5.5, 2.5, 2.5, 1.8])

    # Footer note if data is missing
    if all(c.get('backlinks', 0) == 0 for c in [ctx["your_site_metrics"]] + ctx["competitors_metrics"]):
        _footer(s, "Backlinks data unavailable — activate DataForSEO Backlinks API or paste counts manually in the input form.")


def slide_backlinks_analysis(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Backlinks Analysis", "Where competitors get their links — quality vs. spam patterns")
    _accent_bar(s)
    ins = ctx["insights"].get("backlinks", {})

    paras = [
        ("Overview", {"size": 14, "bold": True, "color": TEAL}),
        (ins.get("narrative", ""), {"size": 11}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0), paras)

    # Two columns: quality backlinks left, spam patterns right
    quality = ins.get("competitor_quality_links", []) or []
    spam = ins.get("spam_patterns", []) or []

    # Left col
    items_q = []
    for q in quality[:3]:
        items_q.append(f"{q.get('domain', '')}: {q.get('comment', '')}")
        for url in (q.get("sample_links", []) or [])[:2]:
            items_q.append(f"   → {url}")
    _bullets(s, Inches(0.5), Inches(3.7), Inches(6.0), Inches(3.4),
             items_q, title="Quality backlink examples", title_color=GREEN, size=9)

    # Right col
    items_s = []
    for sp in spam[:3]:
        items_s.append(f"{sp.get('domain', '')}: {sp.get('pattern', '')}")
        if sp.get("recommendation"):
            items_s.append(f"   → {sp['recommendation']}")
    _bullets(s, Inches(6.8), Inches(3.7), Inches(6.0), Inches(3.4),
             items_s, title="Spam patterns to avoid", title_color=RED, size=9)


def slide_backlink_strategy(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Backlink Building Strategy", "5 methods + monthly PR backlink target")
    _accent_bar(s, color=TEAL)
    ins = ctx["insights"].get("backlinks", {})

    obj = ins.get("objective", "")
    methods = ins.get("methods", {}) or {}
    if obj:
        _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.8),
                        [("Objective", {"size": 13, "bold": True, "color": TEAL}),
                         (obj, {"size": 12})])

    method_blocks = [
        ("Guest Posting", methods.get("guest_posting", ""), TEAL),
        ("Competitor Backlinks", methods.get("competitor_backlinks", ""), AMBER),
        ("Forum Engagement", methods.get("forum", ""), GREEN),
        ("Citation Building", methods.get("citation_building", ""), TEAL),
        ("Local Listings", methods.get("local_listing", ""), AMBER),
    ]
    y = 2.7
    col_w = (12.3 - 0.4) / 5
    for i, (title, text, color) in enumerate(method_blocks):
        x = Inches(0.5 + i * (col_w + 0.1))
        # Title bar
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(y), Inches(col_w), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        tb = s.shapes.add_textbox(x, Inches(y + 0.04), Inches(col_w), Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Content
        _add_paragraphs(s, x, Inches(y + 0.7), Inches(col_w), Inches(3.5),
                        [(text, {"size": 9, "color": NAVY})])


def slide_gmb_overview(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Google My Business — Audit Overview", "Missing elements, strengths, and competitor signals")
    _accent_bar(s)
    ins = ctx["insights"].get("gmb", {})

    # Two columns: missing | plus
    missing = ins.get("missing_elements", []) or []
    plus = ins.get("plus_points", []) or []
    comp = ins.get("competitor_strengths", []) or []

    _bullets(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.7),
             missing[:7], title="Missing on our GMB", title_color=RED, size=11)
    _bullets(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.7),
             plus[:7], title="What we're doing well", title_color=GREEN, size=11)

    _bullets(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.4),
             [f"{c.get('domain', '')}: {c.get('strength', '')}" for c in comp[:6]],
             title="Competitor GMB strengths to learn from", title_color=AMBER, size=10)


def slide_gmb_modifications(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "GMB — Recommended Modifications", "Specific changes by area")
    _accent_bar(s, color=TEAL)
    ins = ctx["insights"].get("gmb", {})
    mods = ins.get("modifications", []) or []

    if not mods:
        _add_paragraphs(s, Inches(0.5), Inches(2), Inches(12.3), Inches(2),
                        [("(No modifications generated.)", {"size": 12})])
        return

    rows = [[m.get("area", ""), m.get("action", "")] for m in mods[:10]]
    _table(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.4 + 0.55 * (len(rows) + 1)),
           headers=["Area", "Recommended Action"],
           rows=rows, col_widths=[3.0, 9.3])


def slide_seo_scorecard(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Technical SEO Health Scorecard", "Pass rate across all audited pages")
    _accent_bar(s)

    seo_checks = ctx["seo_check_rates"]
    labels = [c[0] for c in seo_checks]
    values = [c[1] for c in seo_checks]
    colors = [_score_to_color_hex(int(v)) for v in values]
    img = _bar_chart_image(labels, values, colors, "Pass rate by check (%)")
    s.shapes.add_picture(img, Inches(0.5), Inches(1.4), width=Inches(8.5))

    _bullets(s, Inches(9.3), Inches(1.4), Inches(3.7), Inches(2.5),
             ctx["insights"].get("technical", {}).get("seo_strengths", [])[:4],
             title="Strengths", title_color=GREEN, size=10)
    _bullets(s, Inches(9.3), Inches(4.2), Inches(3.7), Inches(2.8),
             ctx["insights"].get("technical", {}).get("seo_weaknesses", [])[:4],
             title="Weaknesses", title_color=RED, size=10)


def slide_seo_per_page(prs, ctx):
    blank = prs.slide_layouts[6]
    pages_with_issues = ctx.get("pages_with_issues", []) or []
    if not pages_with_issues:
        return
    chunks = [pages_with_issues[i:i + 12] for i in range(0, len(pages_with_issues), 12)]
    for idx, chunk in enumerate(chunks):
        s = prs.slides.add_slide(blank)
        _add_title(s, f"SEO Issues by Page ({idx + 1}/{len(chunks)})",
                   "Specific URLs and their issues")
        _accent_bar(s, color=RED)
        rows = []
        for p in chunk:
            url_short = p["url"]
            if len(url_short) > 75:
                url_short = url_short[:72] + "..."
            issues = "; ".join(p["issues"][:3])
            if len(p["issues"]) > 3:
                issues += f"  (+{len(p['issues']) - 3} more)"
            rows.append([url_short, str(p["score"]), str(len(p["issues"])), issues])
        _table(s, Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.4 + 0.45 * (len(rows) + 1)),
               headers=["URL", "Score", "Issues", "Top Issues"],
               rows=rows, col_widths=[5.0, 0.8, 0.9, 6.0])


def slide_core_web_vitals(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Core Web Vitals — Sampled", "Mobile Lighthouse scores from Google PageSpeed Insights")
    _accent_bar(s)
    psi = ctx.get("pagespeed_results", [])
    if not psi:
        _add_paragraphs(s, Inches(0.5), Inches(2), Inches(12), Inches(1),
                        [("PageSpeed data not available for this run.", {"size": 14})])
        return
    rows = []
    for r in psi:
        url_short = r["url"]
        if len(url_short) > 65:
            url_short = url_short[:62] + "..."
        rows.append([url_short, str(r["performance"]), str(r["seo"]),
                     str(r["lcp_ms"]), f"{r['cls']:.2f}",
                     str(r["inp_ms"]) if r["inp_ms"] else "—"])
    _table(s, Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.4 + 0.5 * (len(rows) + 1)),
           headers=["URL", "Performance", "SEO", "LCP (ms)", "CLS", "INP (ms)"],
           rows=rows, col_widths=[5.0, 1.5, 1.2, 1.7, 1.3, 2.0])


def slide_geo_overview(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "Generative Engine Optimization (GEO)", "Visibility to ChatGPT, Perplexity, Claude, Gemini & AI Overviews")
    _accent_bar(s, color=TEAL)
    geo = ctx["insights"].get("geo", {})

    # KPI row
    geo_score = ctx.get("geo_score", 0)
    citation_avg = ctx.get("citation_avg", 0)
    has_llms = ctx.get("has_llms_txt", False)
    _kpi_card(s, Inches(0.5), Inches(1.5), Inches(3), Inches(1.4),
              "AI-Friendly Score", f"{geo_score}/100", _score_color(geo_score))
    _kpi_card(s, Inches(3.7), Inches(1.5), Inches(3), Inches(1.4),
              "llms.txt", "YES" if has_llms else "NO",
              GREEN if has_llms else RED)
    _kpi_card(s, Inches(6.9), Inches(1.5), Inches(3), Inches(1.4),
              "AI Crawlers Allowed", f"{ctx['ai_allowed']}/{ctx['ai_total']}",
              _score_color(int(100 * ctx['ai_allowed'] / max(1, ctx['ai_total']))))
    _kpi_card(s, Inches(10.1), Inches(1.5), Inches(2.7), Inches(1.4),
              "Citation Signals", f"{citation_avg}/100", _score_color(citation_avg))

    _add_paragraphs(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(3.5),
                    [("Overview", {"size": 14, "bold": True, "color": TEAL}),
                     (geo.get("overview", ""), {"size": 11, "space_before": 4})])


def slide_ai_crawler_table(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "AI Crawler Access", "robots.txt rules for known generative-AI crawlers")
    _accent_bar(s, color=TEAL)

    crawlers = ctx.get("crawlers", [])[:14]
    rows = [[c["name"], c["purpose"], "ALLOWED" if c["allowed"] else "BLOCKED"] for c in crawlers]
    _table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.4 + 0.42 * (len(rows) + 1)),
           headers=["Crawler", "Purpose / Engine", "Access"],
           rows=rows, col_widths=[2.8, 6.7, 2.8],
           status_col=2)


def slide_geo_strategy_text(prs, ctx):
    """GEO narrative slide — llms.txt, citation readiness, AI overview, structured data, entities."""
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    _add_title(s, "GEO Strategy — Detailed Narrative", "All five GEO pillars for your niche")
    _accent_bar(s, color=TEAL)
    geo = ctx["insights"].get("geo", {})
    if not isinstance(geo, dict):
        geo = {}

    paras = [
        ("llms.txt Recommendation", {"size": 12, "bold": True, "color": TEAL}),
        (geo.get("llms_txt_recommendation", ""), {"size": 10}),
        ("Citation Readiness", {"size": 12, "bold": True, "color": TEAL, "space_before": 6}),
        (geo.get("citation_readiness_analysis", ""), {"size": 10}),
        ("AI Overview Strategy", {"size": 12, "bold": True, "color": TEAL, "space_before": 6}),
        (geo.get("ai_overview_strategy", ""), {"size": 10}),
        ("Structured Data Strategy", {"size": 12, "bold": True, "color": TEAL, "space_before": 6}),
        (geo.get("structured_data_strategy", ""), {"size": 10}),
        ("Entity Optimization", {"size": 12, "bold": True, "color": TEAL, "space_before": 6}),
        (geo.get("entity_optimization", ""), {"size": 10}),
    ]
    _add_paragraphs(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.6), paras)


def slide_geo_query_targets(prs, ctx):
    """AI answer-engine queries we want to be cited for."""
    blank = prs.slide_layouts[6]
    geo = ctx["insights"].get("geo", {})
    if not isinstance(geo, dict):
        return
    queries = [q for q in (geo.get("answer_engine_query_targets", []) or []) if isinstance(q, str) and q.strip()]
    formats = [f for f in (geo.get("content_format_recommendations", []) or []) if isinstance(f, dict)]
    if not queries and not formats:
        return
    s = prs.slides.add_slide(blank)
    _add_title(s, "AI Answer-Engine Targets",
               "Queries we should rank as the cited source for + content formats LLMs prefer")
    _accent_bar(s, color=TEAL)

    if queries:
        _bullets(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5),
                 queries[:18], title="Target queries to win citations for", title_color=TEAL, size=10)
    if formats:
        items = [f"{f.get('format', '')} — {f.get('where_to_use', '')}" for f in formats[:8]]
        _bullets(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5),
                 items, title="Content formats LLMs cite most", title_color=AMBER, size=10)


def slide_geo_action_plan(prs, ctx):
    """The 15-item GEO action plan, spread across 2 slides if needed."""
    blank = prs.slide_layouts[6]
    geo = ctx["insights"].get("geo", {})
    if not isinstance(geo, dict):
        return
    actions = [a for a in (geo.get("geo_action_plan", []) or []) if isinstance(a, dict)]
    if not actions:
        return

    chunks = [actions[i:i + 8] for i in range(0, len(actions), 8)]
    for idx, chunk in enumerate(chunks):
        s = prs.slides.add_slide(blank)
        _add_title(s, f"GEO Action Plan ({idx + 1}/{len(chunks)})",
                   "Prioritised tactics to capture AI search citations")
        _accent_bar(s, color=TEAL)
        rows = []
        for a in chunk:
            pri = (a.get("priority", "") or "").upper()
            owner = a.get("owner_role", "") or ""
            effort = a.get("estimated_effort", "") or ""
            rows.append([pri, a.get("action", ""), a.get("why", ""), owner, effort])
        _table(s, Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.4 + 0.55 * (len(rows) + 1)),
               headers=["Priority", "Action", "Why it matters", "Owner", "Effort"],
               rows=rows, col_widths=[1.0, 4.5, 5.0, 1.2, 1.0])


def slide_recommendations(prs, ctx):
    blank = prs.slide_layouts[6]
    recs = ctx["insights"].get("technical", {}).get("priority_recommendations", []) or []
    if not recs:
        return
    s = prs.slides.add_slide(blank)
    _add_title(s, "Priority Recommendations", "Ranked by impact-to-effort (Claude analysis)")
    _accent_bar(s, color=TEAL)

    # Two columns
    col_w = Inches(6.1)
    col_h = Inches(5.6)
    chunks = [recs[:5], recs[5:10]]
    for col_idx, col_recs in enumerate(chunks):
        x = Inches(0.5) + col_idx * Inches(6.4)
        tb = s.shapes.add_textbox(x, Inches(1.4), col_w, col_h)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for r in col_recs:
            priority = r.get("priority", "?")
            title = r.get("title", "")
            why = r.get("why", "")
            action = r.get("action", "")
            area = r.get("area", "")
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(7)
            p.text = f"#{priority}  [{area}]  {title}"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = TEAL

            p2 = tf.add_paragraph()
            p2.text = f"Why: {why}"
            p2.font.size = Pt(9)
            p2.font.color.rgb = NAVY

            p3 = tf.add_paragraph()
            p3.text = f"Action: {action}"
            p3.font.size = Pt(9)
            p3.font.italic = True
            p3.font.color.rgb = GRAY


def slide_closing(prs, ctx):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(12), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Ready to dominate organic & AI search."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(12), Inches(0.8))
    p = tb2.text_frame.paragraphs[0]
    p.text = "Execute the priority recommendations. Re-audit in 30 days."
    p.font.size = Pt(18)
    p.font.color.rgb = TEAL

    tb3 = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(12), Inches(0.4))
    p = tb3.text_frame.paragraphs[0]
    p.text = f"SEO + GEO Strategy & Audit  •  {ctx['domain']}  •  {datetime.now().strftime('%Y-%m-%d')}"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT


# ---------- MASTER BUILDER ----------
def build_report(out_path: str, ctx: dict) -> str:
    """Build the deck. Each slide is wrapped in try/except so a malformed
    Gemini response on ONE section never aborts the whole deck."""
    import traceback
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    errors: list[str] = []

    def safe(name: str, fn, *args):
        try:
            fn(*args)
        except Exception as e:  # noqa: BLE001
            errors.append(f"{name}: {type(e).__name__}: {e}")
            traceback.print_exc()

    safe("cover", slide_cover, prs, ctx)
    safe("scope", slide_scope, prs, ctx)
    safe("toc", slide_toc, prs, ctx)

    safe("dr_explainer", slide_dr_explainer, prs, ctx)
    safe("dr_table", slide_dr_table, prs, ctx)
    safe("dr_analysis", slide_dr_analysis, prs, ctx)

    safe("traffic_table", slide_traffic_table, prs, ctx)
    safe("traffic_analysis", slide_traffic_analysis, prs, ctx)
    safe("competitor_traffic", slide_competitor_traffic_analysis, prs, ctx)

    safe("keyword_homepage", slide_keyword_homepage, prs, ctx)
    safe("keyword_services", slide_keyword_services, prs, ctx)
    safe("keyword_additional", slide_keyword_additional, prs, ctx)

    safe("navigation", slide_navigation, prs, ctx)

    safe("content_intro", slide_content_strategy_intro, prs, ctx)
    silos_raw = ctx["insights"].get("silos", {})
    if not isinstance(silos_raw, dict):
        silos_raw = {}
    silos = [s for s in (silos_raw.get("silos", []) or []) if isinstance(s, dict)]
    for idx, silo in enumerate(silos, start=1):
        safe(f"silo_{idx}", slide_silo, prs, silo, idx, len(silos))

    safe("gmb_overview", slide_gmb_overview, prs, ctx)
    safe("gmb_modifications", slide_gmb_modifications, prs, ctx)

    safe("seo_scorecard", slide_seo_scorecard, prs, ctx)
    safe("seo_per_page", slide_seo_per_page, prs, ctx)
    safe("core_web_vitals", slide_core_web_vitals, prs, ctx)

    safe("geo_overview", slide_geo_overview, prs, ctx)
    safe("ai_crawler", slide_ai_crawler_table, prs, ctx)
    safe("geo_strategy_text", slide_geo_strategy_text, prs, ctx)
    safe("geo_query_targets", slide_geo_query_targets, prs, ctx)
    safe("geo_action_plan", slide_geo_action_plan, prs, ctx)

    safe("recommendations", slide_recommendations, prs, ctx)
    safe("closing", slide_closing, prs, ctx)

    prs.save(out_path)

    # Expose any per-slide errors so the caller can show them in the UI
    ctx["_slide_errors"] = errors
    return out_path
